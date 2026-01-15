"""
PowerPoint Bridge Module
Handles reading and updating PowerPoint presentations using python-pptx.
Works on both Windows and macOS.
"""

import logging
from pathlib import Path
from typing import Any, Optional, Dict, List, Tuple

from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)


class PowerPointBridge:
    """Bridge for reading and updating PowerPoint presentations."""

    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        self.presentation = None
        self._shape_cache = {}
        self._table_cache = {}

    def open(self) -> bool:
        try:
            if not self.file_path.exists():
                logger.error(f"PowerPoint file not found: {self.file_path}")
                return False
            self.presentation = Presentation(str(self.file_path))
            logger.info(f"Opened PowerPoint: {self.file_path}")
            logger.info(f"Slides: {len(self.presentation.slides)}")
            return True
        except Exception as e:
            logger.error(f"Failed to open PowerPoint: {e}")
            return False

    def save(self, output_path: Optional[str] = None) -> bool:
        try:
            save_path = output_path or str(self.file_path)
            self.presentation.save(save_path)
            logger.info(f"Saved PowerPoint to: {save_path}")
            return True
        except Exception as e:
            logger.error(f"Failed to save PowerPoint: {e}")
            return False

    def get_slide_count(self) -> int:
        if self.presentation is None:
            return 0
        return len(self.presentation.slides)

    def _get_slide(self, slide_index: int):
        if self.presentation is None:
            return None
        if slide_index < 1 or slide_index > len(self.presentation.slides):
            logger.error(f"Slide index {slide_index} out of range")
            return None
        return self.presentation.slides[slide_index - 1]

    def _find_shape_by_name(self, slide_index: int, shape_name: str):
        cache_key = (slide_index, shape_name)
        if cache_key in self._shape_cache:
            return self._shape_cache[cache_key]
        slide = self._get_slide(slide_index)
        if slide is None:
            return None
        for shape in slide.shapes:
            if shape.name == shape_name:
                self._shape_cache[cache_key] = shape
                return shape
        logger.warning(f"Shape '{shape_name}' not found on slide {slide_index}")
        return None

    def _find_table_by_name(self, slide_index: int, table_name: str):
        cache_key = (slide_index, table_name)
        if cache_key in self._table_cache:
            return self._table_cache[cache_key]
        slide = self._get_slide(slide_index)
        if slide is None:
            return None
        for shape in slide.shapes:
            if shape.name == table_name and shape.has_table:
                self._table_cache[cache_key] = shape.table
                return shape.table
        logger.warning(f"Table '{table_name}' not found on slide {slide_index}")
        return None

    def update_shape_text(self, slide_index: int, shape_name: str, new_text: str) -> Tuple[bool, str]:
        try:
            shape = self._find_shape_by_name(slide_index, shape_name)
            if shape is None:
                return False, f"Shape '{shape_name}' not found on slide {slide_index}"
            if not shape.has_text_frame:
                return False, f"Shape '{shape_name}' does not contain text"
            text_frame = shape.text_frame
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                if first_para.runs:
                    first_run = first_para.runs[0]
                    font_name = first_run.font.name
                    font_size = first_run.font.size
                    font_bold = first_run.font.bold
                    font_italic = first_run.font.italic
                    font_color = None
                    try:
                        if first_run.font.color and first_run.font.color.rgb:
                            font_color = first_run.font.color.rgb
                    except:
                        pass
                    first_para.clear()
                    run = first_para.add_run()
                    run.text = new_text
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if font_bold is not None:
                        run.font.bold = font_bold
                    if font_italic is not None:
                        run.font.italic = font_italic
                    if font_color:
                        run.font.color.rgb = font_color
                else:
                    first_para.text = new_text
            else:
                text_frame.text = new_text
            return True, f"Updated shape '{shape_name}' on slide {slide_index}"
        except Exception as e:
            logger.error(f"Error updating shape: {e}")
            return False, str(e)

    def update_table_cell(self, slide_index: int, table_name: str, row: int, col: int, new_text: str) -> Tuple[bool, str]:
        try:
            table = self._find_table_by_name(slide_index, table_name)
            if table is None:
                return False, f"Table '{table_name}' not found on slide {slide_index}"
            row_idx, col_idx = row - 1, col - 1
            if row_idx < 0 or row_idx >= len(table.rows):
                return False, f"Row {row} out of range"
            if col_idx < 0 or col_idx >= len(table.columns):
                return False, f"Column {col} out of range"
            cell = table.cell(row_idx, col_idx)
            if cell.text_frame.paragraphs:
                first_para = cell.text_frame.paragraphs[0]
                if first_para.runs:
                    first_run = first_para.runs[0]
                    font_name = first_run.font.name
                    font_size = first_run.font.size
                    font_bold = first_run.font.bold
                    font_color = None
                    try:
                        if first_run.font.color and first_run.font.color.rgb:
                            font_color = first_run.font.color.rgb
                    except:
                        pass
                    first_para.clear()
                    run = first_para.add_run()
                    run.text = new_text
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if font_bold is not None:
                        run.font.bold = font_bold
                    if font_color:
                        run.font.color.rgb = font_color
                else:
                    first_para.text = new_text
            else:
                cell.text = new_text
            return True, f"Updated cell ({row},{col}) in '{table_name}'"
        except Exception as e:
            return False, str(e)

    def list_shapes(self, slide_index: int) -> List[Dict[str, Any]]:
        shapes_info = []
        slide = self._get_slide(slide_index)
        if slide is None:
            return shapes_info
        for shape in slide.shapes:
            info = {
                'name': shape.name,
                'type': str(shape.shape_type),
                'has_text': shape.has_text_frame,
                'has_table': shape.has_table
            }
            if shape.has_text_frame:
                info['text_preview'] = shape.text_frame.text[:50] if shape.text_frame.text else ''
            if shape.has_table:
                info['table_size'] = f"{len(shape.table.rows)}x{len(shape.table.columns)}"
            shapes_info.append(info)
        return shapes_info

    def list_tables(self, slide_index: int) -> List[Dict[str, Any]]:
        tables_info = []
        slide = self._get_slide(slide_index)
        if slide is None:
            return tables_info
        for shape in slide.shapes:
            if shape.has_table:
                tables_info.append({
                    'name': shape.name,
                    'rows': len(shape.table.rows),
                    'columns': len(shape.table.columns)
                })
        return tables_info


def check_presentation(file_path: str) -> Tuple[bool, str]:
    bridge = PowerPointBridge(file_path)
    if bridge.open():
        return True, f"OK: {bridge.file_path.name} ({bridge.get_slide_count()} slides)"
    return False, f"Failed to open {file_path}"
