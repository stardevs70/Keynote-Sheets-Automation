#!/usr/bin/env python3
"""
Create a sample PowerPoint file for testing the automation.
This creates a presentation with named shapes and tables that can be updated.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

def create_sample_presentation(output_path: str = None):
    """Create a sample PowerPoint presentation with named shapes and tables."""
    if output_path is None:
        # Default to same directory as this script
        script_dir = Path(__file__).parent
        output_path = str(script_dir / "sample_investor_deck.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # Slide 1: Title Slide
    # =========================================================================
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(blank_layout)

    # Title shape
    title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    title.name = "Title"
    tf = title.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Investor Report"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Date shape
    date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.5))
    date_box.name = "ReportDate"
    tf = date_box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Q4 2024"
    run.font.size = Pt(24)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # =========================================================================
    # Slide 2: Key Metrics
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)

    # Section title
    section_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    section_title.name = "SectionTitle"
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Performance Metrics"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Revenue shape
    revenue_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(3), Inches(1.5))
    revenue_box.name = "RevenueBox"
    revenue_box.fill.solid()
    revenue_box.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0x99)
    tf = revenue_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Revenue"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "$5,000,000"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    # Create separate shape for revenue value that can be updated
    revenue_value = slide2.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(3), Inches(0.6))
    revenue_value.name = "RevenueValue"
    tf = revenue_value.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "$5,000,000"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)

    # Growth Rate shape
    growth_box = slide2.shapes.add_textbox(Inches(4), Inches(1.2), Inches(3), Inches(1.5))
    growth_box.name = "GrowthRate"
    tf = growth_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Growth Rate"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = "Arial"
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "25.5%"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x99, 0x33)

    # Growth Rate Value shape (separate for easy updates)
    growth_value = slide2.shapes.add_textbox(Inches(4), Inches(2.0), Inches(3), Inches(0.6))
    growth_value.name = "GrowthValue"
    tf = growth_value.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "25.5%"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x99, 0x33)

    # Customer Count shape
    customers_box = slide2.shapes.add_textbox(Inches(7.5), Inches(1.2), Inches(3), Inches(1.5))
    customers_box.name = "CustomerCount"
    tf = customers_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Active Customers"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = "Arial"
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "1,250"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Customer Value shape (separate for updates)
    customer_value = slide2.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(3), Inches(0.6))
    customer_value.name = "CustomerValue"
    tf = customer_value.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "1,250"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # =========================================================================
    # Slide 3: Financial Table
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)

    # Section title
    section_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    section_title.name = "FinancialTitle"
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Financial Summary"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Create a table
    rows, cols = 5, 4
    table_shape = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(10), Inches(3))
    table_shape.name = "FinancialTable"
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2.333)
    table.columns[2].width = Inches(2.333)
    table.columns[3].width = Inches(2.333)

    # Header row
    headers = ["Metric", "Q2 2024", "Q3 2024", "Q4 2024"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        para.font.name = "Arial"

    # Data rows
    data = [
        ["Revenue", "$4,200,000", "$4,600,000", "$5,000,000"],
        ["Expenses", "$2,800,000", "$3,000,000", "$3,200,000"],
        ["Net Income", "$1,400,000", "$1,600,000", "$1,800,000"],
        ["EBITDA", "$1,800,000", "$2,000,000", "$2,200,000"],
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_value
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.name = "Arial"
            if col_idx == 0:
                para.font.bold = True
                para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            else:
                para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # =========================================================================
    # Slide 4: KPI Table
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)

    # Section title
    section_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    section_title.name = "KPITitle"
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Performance Indicators"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # KPI Table
    rows, cols = 6, 3
    kpi_table_shape = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(8), Inches(3.5))
    kpi_table_shape.name = "KPITable"
    kpi_table = kpi_table_shape.table

    # Set column widths
    kpi_table.columns[0].width = Inches(4)
    kpi_table.columns[1].width = Inches(2)
    kpi_table.columns[2].width = Inches(2)

    # Header row
    kpi_headers = ["KPI", "Current", "Target"]
    for col, header in enumerate(kpi_headers):
        cell = kpi_table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0x99)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        para.font.name = "Arial"

    # KPI Data
    kpi_data = [
        ["Customer Acquisition Cost", "$150", "$125"],
        ["Customer Lifetime Value", "$2,500", "$3,000"],
        ["Churn Rate", "5.2%", "4.0%"],
        ["Net Promoter Score", "72", "80"],
        ["Monthly Active Users", "45,000", "50,000"],
    ]

    for row_idx, row_data in enumerate(kpi_data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = kpi_table.cell(row_idx, col_idx)
            cell.text = cell_value
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.name = "Arial"
            if col_idx == 0:
                para.font.bold = True

    # =========================================================================
    # Slide 5: Summary with Text Shapes
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)

    # Section title
    section_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    section_title.name = "SummaryTitle"
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Executive Summary"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Period text
    period_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
    period_box.name = "ReportPeriod"
    tf = period_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Reporting Period: Q4 2024"
    run.font.size = Pt(16)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Total Revenue box
    total_revenue = slide5.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4), Inches(0.8))
    total_revenue.name = "TotalRevenue"
    tf = total_revenue.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Total Revenue: $5,000,000"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # YoY Growth box
    yoy_growth = slide5.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4), Inches(0.8))
    yoy_growth.name = "YoYGrowth"
    tf = yoy_growth.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Year-over-Year Growth: 25.5%"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x99, 0x33)

    # Customer count summary
    customer_summary = slide5.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(4), Inches(0.8))
    customer_summary.name = "CustomerSummary"
    tf = customer_summary.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Total Customers: 1,250"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Inception date
    inception_date = slide5.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5), Inches(0.5))
    inception_date.name = "InceptionDate"
    tf = inception_date.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Fund Inception: April 1, 2021"
    run.font.size = Pt(14)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Save the presentation
    output = Path(output_path)
    prs.save(str(output))
    print(f"Created sample presentation: {output.absolute()}")

    # Print shape inventory for reference
    print("\n" + "=" * 60)
    print("SHAPE INVENTORY FOR TESTING")
    print("=" * 60)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"\nSlide {slide_idx}:")
        for shape in slide.shapes:
            shape_type = "TABLE" if shape.has_table else "TEXT" if shape.has_text_frame else "OTHER"
            text_preview = ""
            if shape.has_text_frame and shape.text_frame.text:
                text_preview = f' -> "{shape.text_frame.text[:30]}..."' if len(shape.text_frame.text) > 30 else f' -> "{shape.text_frame.text}"'
            if shape.has_table:
                text_preview = f' ({len(shape.table.rows)}x{len(shape.table.columns)} table)'
            print(f"  - {shape.name} [{shape_type}]{text_preview}")

    return str(output.absolute())


if __name__ == '__main__':
    create_sample_presentation()
